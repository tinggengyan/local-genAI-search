import PyPDF2
from os import listdir
from os.path import isfile, join,isdir

import torch
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_qdrant import Qdrant
import sys
from langchain_text_splitters import TokenTextSplitter
from pptx import Presentation
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams
import docx
import os
from environment_var import qdrant_url, qdrant_api_key, model_name, collection_name

def get_files(dir):
    file_list = []
    for dir, _, filenames in os.walk(dir):
        for f in filenames:
            file_list.append(os.path.join(dir, f))
    return file_list

def getTextFromWord(filename):
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)

def getTextFromPPTX(filename):
    prs = Presentation(filename)
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:  # 只添加非空文本
                        fullText.append(text)
            except AttributeError:
                # 跳过没有 text 属性的形状（如图片）
                continue
            except Exception as e:
                print(f"Warning: Error processing shape in {filename}: {str(e)}")
                continue
    return '\n'.join(fullText)

def main_indexing(mypath):
    # 使用 environment_var.py 中定义的模型
    if torch.cuda.is_available():
        model_kwargs = {'device': 'cuda'}
    elif torch.backends.mps.is_available():
        model_kwargs = {'device': 'mps'}
    else:
        model_kwargs = {'device': 'cpu'}
    encode_kwargs = {'normalize_embeddings': True}
    hf = HuggingFaceEmbeddings(
        model_name=model_name,
        model_kwargs=model_kwargs,
        encode_kwargs=encode_kwargs
    )
    client = QdrantClient(url=qdrant_url, api_key=qdrant_api_key)
    if client.collection_exists(collection_name):
        client.delete_collection(collection_name)

    # 使用 512 维向量，因为 bge-small-zh-v1.5 模型生成的是 512 维向量
    client.create_collection(collection_name,vectors_config=VectorParams(size=512, distance=Distance.DOT))
    qdrant = Qdrant(client, collection_name, hf)
    print("Indexing...")
    onlyfiles = get_files(mypath)
    for file in onlyfiles:
        try:
            file_content = ""
            if file.find("~") > 0:  # 跳过临时文件
                print(f"Skipping temporary file: {file}")
                continue
            elif file.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp')):
                print(f"Skipping image file: {file}")
                continue
            elif file.endswith(".pdf"):
                try:
                    print("indexing "+file)
                    reader = PyPDF2.PdfReader(file)
                    for i in range(0,len(reader.pages)):
                        file_content = file_content + " "+reader.pages[i].extract_text()
                except Exception as exc:
                    print(f"Error processing PDF {file}: {str(exc)}")
                    continue
            elif file.endswith(".txt") or file.endswith(".md") or file.endswith(".markdown"):
                try:
                    print("indexing " + file)
                    f = open(file,'r',encoding='utf-8',errors='ignore')
                    file_content = f.read()
                    f.close()
                except Exception as exc:
                    print(f"Error processing text file {file}: {str(exc)}")
                    continue
            elif file.endswith(".docx"):
                try:
                    print("indexing " + file)
                    file_content = getTextFromWord(file)
                except Exception as exc:
                    print(f"Error processing Word file {file}: {str(exc)}")
                    continue
            elif file.endswith(".pptx"):
                try:
                    print("indexing " + file)
                    file_content = getTextFromPPTX(file)
                except Exception as exc:
                    print(f"Error processing PowerPoint file {file}: {str(exc)}")
                    continue
            else:
                print(f"Skipping unsupported file type: {file}")
                continue

            if not file_content.strip():
                print(f"Warning: No content extracted from {file}")
                continue

            text_splitter = TokenTextSplitter(chunk_size=500, chunk_overlap=50)
            texts = text_splitter.split_text(file_content)
            if not texts:
                print(f"Warning: No text chunks created for {file}")
                continue

            metadata = []
            for i in range(0,len(texts)):
                metadata.append({"path":file})
            qdrant.add_texts(texts,metadatas=metadata)
            print(f"Successfully indexed {file} with {len(texts)} chunks")
        except Exception as e:
            print(f"Error processing file {file}: {str(e)}")
            continue

    print("Finished indexing!")

if __name__ == "__main__":
    arguments = sys.argv
    if len(arguments)>1:
        main_indexing(arguments[1])
    else:
        print("You need to provide a path to folder with documents to index as command line argument")